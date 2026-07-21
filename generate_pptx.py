"""
generate_pptx.py

Assembles the PNG charts produced by notebook/analysis.ipynb into a
client-ready PowerPoint deck (data/client_presentation.pptx).

Run after executing the notebook:
    python generate_pptx.py
"""

import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEAL = RGBColor(0x0F, 0x4C, 0x45)
CORAL = RGBColor(0xE0, 0x65, 0x4A)
INK = RGBColor(0x16, 0x30, 0x2B)

CHART_DIR = "data/charts"
OUTPUT_PATH = "data/client_presentation.pptx"

SLIDES = [
    {
        "title": "Ages With the Highest Income",
        "subtitle": "Average total income by respondent age",
        "image": "income_by_age.png",
    },
    {
        "title": "Spending by Category, by Gender",
        "subtitle": "Average amount spent per expense category, broken out by gender",
        "image": "spending_by_gender.png",
    },
    {
        "title": "Income vs. Expenses by Age",
        "subtitle": "Disposable income context for the healthcare product launch",
        "image": "income_vs_expenses_by_age.png",
    },
]


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = TEAL

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.2))
    tf = tb.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    sb = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.7), Inches(0.8))
    sf = sb.text_frame
    sf.text = subtitle
    sf.paragraphs[0].font.size = Pt(18)
    sf.paragraphs[0].font.color.rgb = RGBColor(0xE7, 0xEE, 0xEC)


def add_chart_slide(prs, title, subtitle, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = TEAL

    sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.95), Inches(12.1), Inches(0.5))
    sf = sub_box.text_frame
    sf.text = subtitle
    sf.paragraphs[0].font.size = Pt(14)
    sf.paragraphs[0].font.color.rgb = RGBColor(0x4A, 0x6B, 0x66)

    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(1.4), Inches(1.65), width=Inches(10.5))
    else:
        placeholder = slide.shapes.add_textbox(Inches(1.4), Inches(3), Inches(10), Inches(1))
        placeholder.text_frame.text = f"[Chart not found: {image_path}]"
        placeholder.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xB4, 0x33, 0x1A)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Household Income & Spending Survey",
        "Findings ahead of the new healthcare product launch",
    )

    for slide_info in SLIDES:
        image_path = os.path.join(CHART_DIR, slide_info["image"])
        add_chart_slide(prs, slide_info["title"], slide_info["subtitle"], image_path)

    os.makedirs(os.path.dirname(OUTPUT_PATH) or ".", exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved presentation to {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
